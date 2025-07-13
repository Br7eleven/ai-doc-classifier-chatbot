import os
import joblib
import PyPDF2
import pandas as pd
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import numpy as np
from docx import Document

# DOCX support
from docx import Document

# PPTX support
from pptx import Presentation

# Excel support
import openpyxl
import xlrd

# RTF support
from striprtf.striprtf import rtf_to_text

# ODT support
from odf.opendocument import load
from odf.text import P, H


app = Flask(__name__, template_folder="templates", static_folder="static")

# Load model and vectorizer
try:
    model = joblib.load("model/classifier_model.joblib")
    vectorizer = joblib.load("model/vectorizer.joblib")
except Exception as e:
    print("[ERROR] Failed to load model/vectorizer:", e)
    model = None
    vectorizer = None

def extract_text(file):
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        return " ".join([page.extract_text() or "" for page in reader.pages])

    elif filename.endswith(".docx"):
        doc = Document(file)
        return " ".join([para.text for para in doc.paragraphs])

    elif filename.endswith(".txt"):
        return file.read().decode("utf-8")

    elif filename.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text

    elif filename.endswith(".xlsx"):
        df = pd.read_excel(file)
        return " ".join(df.astype(str).values.flatten())

    else:
        raise ValueError("Unsupported file type")

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/upload", methods=["POST"])
def upload():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    text = extract_text(file)

    if not text.strip():
        return jsonify({"error": "File is empty or unsupported format"}), 400

    if model and vectorizer:
        try:
            features = vectorizer.transform([text])
            prediction = model.predict(features)[0]
            confidence = np.max(model.predict_proba(features))
            return jsonify({"category": prediction, "confidence": round(float(confidence), 2)})
        except Exception as e:
            return jsonify({"error": str(e)}), 500
    else:
        return jsonify({"error": "Model not loaded"}), 500

if __name__ == "__main__":
    app.run(port=5001, debug=True)
